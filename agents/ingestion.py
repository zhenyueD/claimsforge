"""
SOP/document ingestion pipeline.

Pipeline:
  1. Parse — PDF / DOCX / Markdown / plain text → raw text
  2. Chunk — semantic chunks (~600 tokens, prefer ### / ## section breaks)
  3. Synthesize — Gemini converts each chunk into a typed KBEntry
                  (HUMAN_SOP source, RULE/PRINCIPLE/CASE type)
  4. Upsert + embed — into the unified KB

After this runs, every agent (Compensation, Needs, Verifier, etc.) automatically
sees the new knowledge through hybrid_search.
"""
from __future__ import annotations

import io
import logging
import re
import time
from pathlib import Path
from typing import Optional

from pydantic import BaseModel, Field

from gemini_client import GeminiError, structured
from unified_kb import KBEntry, KBSource, KBType, make_id, upsert
from embedding_index import index_all

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────
#  1) Parse
# ─────────────────────────────────────────────────────────
def parse_pdf(blob: bytes) -> str:
    import pypdf
    reader = pypdf.PdfReader(io.BytesIO(blob))
    return "\n\n".join((p.extract_text() or "") for p in reader.pages).strip()


def parse_docx(blob: bytes) -> str:
    import docx
    d = docx.Document(io.BytesIO(blob))
    paras = [p.text for p in d.paragraphs if p.text.strip()]
    # also pull tables
    for tbl in d.tables:
        for row in tbl.rows:
            paras.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(paras).strip()


def parse_pptx(blob: bytes) -> str:
    """Extract text from every slide in a .pptx deck."""
    import pptx
    prs = pptx.Presentation(io.BytesIO(blob))
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = [f"## Slide {i}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                parts.append(shape.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        out.append("\n".join(parts))
    return "\n\n".join(out).strip()


def parse_xlsx(blob: bytes) -> str:
    """Extract cells from every sheet of an .xlsx workbook."""
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(blob), data_only=True, read_only=True)
    out = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            out.append(f"## Sheet: {ws.title}\n" + "\n".join(rows))
    return "\n\n".join(out).strip()


def parse_xls(blob: bytes) -> str:
    """Legacy .xls — xlrd only handles up to .xls (2003)."""
    import xlrd
    wb = xlrd.open_workbook(file_contents=blob)
    out = []
    for sheet in wb.sheets():
        rows = []
        for r in range(sheet.nrows):
            cells = [str(sheet.cell_value(r, c)).strip() for c in range(sheet.ncols)]
            cells = [c for c in cells if c]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            out.append(f"## Sheet: {sheet.name}\n" + "\n".join(rows))
    return "\n\n".join(out).strip()


def parse_via_gemini_file_api(blob: bytes, filename: str) -> str:
    """Fallback: upload the bytes to Gemini and ask it to transcribe.
    Handles .doc, .ppt, and anything else we can't parse locally."""
    from google import genai as _genai
    from google.genai import types as _gtypes
    import gemini_client as _gc
    client = _gc.get_client()

    # Upload as a temp file (Gemini file API)
    import tempfile
    suffix = Path(filename).suffix or ".bin"
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(blob)
        tmp_path = tmp.name

    try:
        uploaded = client.files.upload(file=tmp_path)
        # Wait for processing
        for _ in range(20):
            f = client.files.get(name=uploaded.name)
            if str(f.state) in ("ACTIVE", "FileState.ACTIVE"):
                break
            time.sleep(0.5)

        resp = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=[
                "Transcribe ALL text content from this document, faithfully. "
                "Preserve section structure with ## headers where appropriate. "
                "Do not summarize. Output plain text only.",
                uploaded,
            ],
            config=_gtypes.GenerateContentConfig(
                temperature=0,
                max_output_tokens=16000,
                thinking_config=_gtypes.ThinkingConfig(thinking_budget=0),
            ),
        )
        text = (resp.text or "").strip()
        try:
            client.files.delete(name=uploaded.name)
        except Exception:
            pass
        return text
    finally:
        try:
            Path(tmp_path).unlink()
        except Exception:
            pass


def parse_text(blob: bytes) -> str:
    for enc in ("utf-8", "utf-16", "gbk", "latin-1"):
        try:
            return blob.decode(enc).strip()
        except UnicodeDecodeError:
            continue
    return blob.decode("utf-8", errors="replace").strip()


def parse_document(filename: str, blob: bytes) -> str:
    lower = filename.lower()
    try:
        if lower.endswith(".pdf"):
            return parse_pdf(blob)
        if lower.endswith(".docx"):
            return parse_docx(blob)
        if lower.endswith(".pptx"):
            return parse_pptx(blob)
        if lower.endswith(".xlsx"):
            return parse_xlsx(blob)
        if lower.endswith(".xls"):
            return parse_xls(blob)
        # Legacy formats: Gemini File API can parse them directly
        if lower.endswith(".doc") or lower.endswith(".ppt"):
            logger.info("legacy format %s → Gemini File API fallback", filename)
            return parse_via_gemini_file_api(blob, filename)
        # Plain text / markdown / json / csv
        return parse_text(blob)
    except Exception as e:
        logger.warning("local parse of %s failed (%s), trying Gemini File API fallback", filename, e)
        try:
            return parse_via_gemini_file_api(blob, filename)
        except Exception as e2:
            logger.error("Gemini fallback also failed for %s: %s", filename, e2)
            raise


# ─────────────────────────────────────────────────────────
#  2) Chunk — semantic-aware, prefer markdown headers
# ─────────────────────────────────────────────────────────
_HEADER_RE = re.compile(r"^(#{1,6}\s|第[一二三四五六七八九十0-9]+[章节条]\b|\d+\.\s|[A-Z][A-Z\s]{3,}:)", re.MULTILINE)


def chunk_text(text: str, *, max_chars: int = 2400, prefer_headers: bool = True) -> list[str]:
    """Chunk text into semantic units. Always splits at markdown / numbered
    headers when present (so multi-section docs produce multi-entry KB);
    falls back to paragraph + length cap."""
    text = text.strip()
    if not text:
        return []

    chunks: list[str] = []

    # If there are 2+ headers, ALWAYS split by header (even if total length is small).
    header_hits = list(_HEADER_RE.finditer(text))
    has_multi_section = len(header_hits) >= 2

    if not has_multi_section and len(text) <= max_chars:
        return [text]

    if prefer_headers and header_hits:
        # Split at headers: each header starts a new chunk.
        pieces = []
        last = 0
        for m in header_hits:
            if m.start() > last:
                pieces.append(text[last:m.start()].strip())
            last = m.start()
        pieces.append(text[last:].strip())

        # Drop the leading preamble if it's just a document title with no body
        if pieces and len(pieces[0]) < 80 and not pieces[0].startswith(("#", "第", "1.", "2.")):
            pieces = pieces[1:]

        # Each non-empty piece becomes its own chunk, but split pieces that exceed max_chars
        for p in pieces:
            if not p:
                continue
            if len(p) <= max_chars:
                chunks.append(p)
            else:
                # paragraph-cap inside an oversize piece
                sub = ""
                for para in p.split("\n\n"):
                    if len(sub) + len(para) + 2 <= max_chars:
                        sub = (sub + "\n\n" + para).strip()
                    else:
                        if sub:
                            chunks.append(sub)
                        sub = para
                if sub:
                    chunks.append(sub)
    else:
        # Plain paragraph split
        para = ""
        for p in text.split("\n\n"):
            p = p.strip()
            if not p:
                continue
            if len(para) + len(p) + 2 <= max_chars:
                para = (para + "\n\n" + p).strip()
            else:
                if para:
                    chunks.append(para)
                para = p
        if para:
            chunks.append(para)

    return [c for c in chunks if c]


# ─────────────────────────────────────────────────────────
#  3) Synthesize — Gemini → typed KBEntry
# ─────────────────────────────────────────────────────────
class SynthesizedChunk(BaseModel):
    title: str = Field(description="≤ 80 chars summary headline")
    domain: str = Field(description="damage / emotion / logistics / pricing / fraud / refund / brand_voice / etc")
    type: KBType
    scenario: str = Field(description="When does this apply (1-2 sentences)")
    decision: str = Field(description="What to do / what the rule says (1-3 sentences)")
    rationale: str = Field(description="Why — the underlying reasoning. Keep concise.")
    tags: list[str] = Field(default_factory=list, description="3-8 short tags")
    customer_facing_name: Optional[str] = Field(
        default=None,
        description="Friendly name to cite to customers (e.g. '30-day return policy'). null if not customer-facing."
    )
    quality_score: float = Field(
        default=0.7, ge=0, le=1,
        description="How directly actionable is this chunk? 0.7 default for human SOPs."
    )


_SYNTH_SYSTEM = """You ingest one chunk of an internal customer-service SOP / policy /
brand-voice doc and convert it into ONE structured KB entry that the multi-agent
system can later retrieve.

Rules:
  - Be faithful to the source. Do NOT invent details.
  - title: a useful 80-char headline (not a generic "Section 3")
  - domain: pick ONE: damage / refund / shipping / return / payment / fraud /
            customer_loyalty / brand_voice / escalation / communication /
            emotion / pricing / policy / inventory / other
  - type: pick ONE — RULE (if-then), PRINCIPLE (general guidance),
          CASE (specific example), DECISION_LOG (what was done)
  - scenario: when does this rule trigger? (factual extraction, no embellishment)
  - decision: what to do? (or what was done)
  - rationale: why? (extract if present; otherwise short inference)
  - customer_facing_name: ONLY if this is a named policy customers should be told
                          about. Otherwise null.
  - tags: 3-8 short keywords for retrieval. Use the source's own language.
  - quality_score: 0.7 for clear actionable chunks; 0.5 for ambiguous; 0.3 for
                   chunks that are mostly headers/TOC noise.

Output language: keep the language of the source (Chinese in, Chinese out).
"""


def synthesize_chunk(chunk_text: str) -> Optional[SynthesizedChunk]:
    try:
        return structured(
            prompt=f"## Source chunk\n\n{chunk_text}",
            schema=SynthesizedChunk,
            system=_SYNTH_SYSTEM,
            temperature=0.2,
            max_tokens=800,
        )
    except GeminiError as e:
        logger.warning("synthesize failed on chunk: %s", e)
        return None


# ─────────────────────────────────────────────────────────
#  4) Ingest end-to-end
# ─────────────────────────────────────────────────────────
class IngestionReport(BaseModel):
    filename: str
    bytes: int
    raw_chars: int
    chunks: int
    entries_written: int
    entry_ids: list[str]
    skipped: int
    errors: list[str] = Field(default_factory=list)
    embedded: int
    elapsed_s: float


def ingest_document(
    filename: str,
    blob: bytes,
    *,
    contributor: str = "human_upload",
    domain_hint: Optional[str] = None,
    max_chunks: Optional[int] = None,
    rebuild_embeddings: bool = True,
) -> IngestionReport:
    t0 = time.monotonic()
    raw = parse_document(filename, blob)
    chunks = chunk_text(raw)
    if max_chunks:
        chunks = chunks[:max_chunks]

    report = IngestionReport(
        filename=filename, bytes=len(blob), raw_chars=len(raw),
        chunks=len(chunks), entries_written=0, entry_ids=[], skipped=0,
        embedded=0, elapsed_s=0,
    )

    for idx, chunk in enumerate(chunks):
        synth = synthesize_chunk(chunk)
        if not synth or synth.quality_score < 0.25:
            report.skipped += 1
            continue
        entry_id = make_id(f"{filename}-c{idx}-{chunk[:80]}")
        entry = KBEntry(
            id=entry_id,
            source=KBSource.HUMAN_SOP,
            type=synth.type,
            domain=domain_hint or synth.domain or "general",
            title=synth.title,
            customer_facing_name=synth.customer_facing_name,
            scenario=synth.scenario,
            decision=synth.decision,
            rationale=synth.rationale,
            tags=synth.tags,
            contributor=contributor,
            quality_score=synth.quality_score,
            source_doc=filename,
            source_chunk=idx,
        )
        upsert(entry)
        report.entry_ids.append(entry_id)
        report.entries_written += 1

    # Build embeddings for the new entries
    if rebuild_embeddings and report.entries_written > 0:
        try:
            idx_result = index_all(rate_limit_sleep=0.05)
            report.embedded = idx_result.get("indexed", 0)
        except Exception as e:
            report.errors.append(f"embed failed: {e}")

    report.elapsed_s = round(time.monotonic() - t0, 2)
    return report
